"""PowerPoint形式のレポート生成ユーティリティ"""
from typing import Dict, Any, List
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PPTXGenerator:
    """PowerPoint形式のレポートを生成"""
    
    def __init__(self):
        """初期化"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # カラーパレット
        self.colors = {
            'primary': RGBColor(0, 102, 204),      # ブルー
            'secondary': RGBColor(255, 153, 0),    # オレンジ
            'success': RGBColor(76, 175, 80),      # グリーン
            'text': RGBColor(33, 33, 33),          # ダークグレー
            'light': RGBColor(245, 245, 245)       # ライトグレー
        }
    
    def generate_report(
        self,
        project_info: Dict[str, Any],
        agent_results: List[Dict[str, Any]],
        output_path: str
    ) -> str:
        """
        PowerPointレポートを生成
        
        Args:
            project_info: プロジェクト情報
            agent_results: 各エージェントの分析結果
            output_path: 出力ファイルパス
        
        Returns:
            生成されたファイルパス
        """
        # タイトルスライド
        self._add_title_slide(project_info)
        
        # エグゼクティブサマリー
        self._add_executive_summary(agent_results)
        
        # 各分析結果
        for result in agent_results:
            if result.get('status') == 'success':
                self._add_analysis_slide(result)
        
        # 推奨事項
        self._add_recommendations_slide(agent_results)
        
        # アクションプラン
        self._add_action_plan_slide()
        
        # 保存
        self.prs.save(output_path)
        return output_path
    
    def _add_title_slide(self, project_info: Dict[str, Any]):
        """タイトルスライドを追加"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白レイアウト
        
        # 背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['primary']
        
        # タイトル
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "戦略コンサルティング統合レポート"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # サブタイトル
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(8), Inches(1.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = (
            f"{project_info.get('client_name', 'N/A')}\n"
            f"{project_info.get('industry', 'N/A')} | "
            f"{datetime.now().strftime('%Y年%m月%d日')}"
        )
        for para in subtitle_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER
    
    def _add_executive_summary(self, agent_results: List[Dict[str, Any]]):
        """エグゼクティブサマリースライドを追加"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # タイトル+コンテンツ
        
        # タイトル
        title = slide.shapes.title
        title.text = "エグゼクティブサマリー"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # コンテンツ
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        successful_analyses = [r for r in agent_results if r.get('status') == 'success']
        
        # 各分析のサマリー
        for result in successful_analyses:
            analysis_type = result.get('analysis_type', 'unknown')
            
            p = text_frame.add_paragraph()
            p.level = 0
            
            if analysis_type == 'market':
                p.text = "📊 市場分析"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.colors['primary']
                
                market_result = result.get('result', {})
                details = [
                    f"市場魅力度: {market_result.get('market_attractiveness', 'N/A')}",
                    f"市場規模: {market_result.get('market_size', 0):,.0f}円",
                    f"成長率: {market_result.get('growth_rate', 0)}%"
                ]
                
            elif analysis_type == 'financial':
                p.text = "💰 財務分析"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.colors['success']
                
                financial_result = result.get('result', {})
                prof = financial_result.get('profitability_ratios', {})
                details = [
                    f"総合評価: {financial_result.get('overall_assessment', 'N/A')}",
                    f"営業利益率: {prof.get('operating_margin', 0):.1f}%"
                ]
                
            elif analysis_type == 'strategy':
                p.text = "🎯 戦略分析"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.colors['secondary']
                
                details = ["複数の戦略フレームワークを用いた包括的分析を実施"]
            
            else:
                continue
            
            # 詳細を追加
            for detail in details:
                p = text_frame.add_paragraph()
                p.text = f"  • {detail}"
                p.level = 1
                p.font.size = Pt(16)
                p.font.color.rgb = self.colors['text']
            
            # 空行
            text_frame.add_paragraph()
    
    def _add_analysis_slide(self, result: Dict[str, Any]):
        """分析結果スライドを追加"""
        analysis_type = result.get('analysis_type', 'unknown')
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # タイトル
        title = slide.shapes.title
        if analysis_type == 'market':
            title.text = "市場分析結果"
            color = self.colors['primary']
        elif analysis_type == 'financial':
            title.text = "財務分析結果"
            color = self.colors['success']
        elif analysis_type == 'strategy':
            title.text = "戦略分析結果"
            color = self.colors['secondary']
        else:
            title.text = "分析結果"
            color = self.colors['text']
        
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = color
        
        # コンテンツ（簡略版）
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # フォーマット済み出力から主要な情報を抽出
        formatted_output = result.get('formatted_output', '')
        
        # 最初の500文字程度を表示（スライドに収まるように）
        summary_text = formatted_output[:500] + "..." if len(formatted_output) > 500 else formatted_output
        
        p = text_frame.add_paragraph()
        p.text = "詳細な分析結果はマークダウンレポートをご参照ください。"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = self.colors['text']
    
    def _add_recommendations_slide(self, agent_results: List[Dict[str, Any]]):
        """推奨事項スライドを追加"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # タイトル
        title = slide.shapes.title
        title.text = "統合的な推奨事項"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # コンテンツ
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        all_recommendations = []
        
        # 各エージェントからの推奨事項を収集
        for result in agent_results:
            if result.get('status') != 'success':
                continue
            
            analysis_result = result.get('result', {})
            
            if result.get('analysis_type') == 'market':
                recs = analysis_result.get('recommendations', [])
                all_recommendations.extend([('市場', rec) for rec in recs])
            elif result.get('analysis_type') == 'financial':
                recs = analysis_result.get('recommendations', [])
                all_recommendations.extend([('財務', rec) for rec in recs])
        
        # 推奨事項を表示
        for i, (category, rec) in enumerate(all_recommendations[:5], 1):  # 最大5件
            p = text_frame.add_paragraph()
            p.text = f"{i}. [{category}] {rec}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text']
            p.space_after = Pt(12)
    
    def _add_action_plan_slide(self):
        """アクションプランスライドを追加"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # タイトル
        title = slide.shapes.title
        title.text = "アクションプラン"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # コンテンツ
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        action_plans = [
            ("短期（1-3ヶ月）", [
                "データ収集と詳細分析の実施",
                "優先度の高い施策の計画立案",
                "ステークホルダーとの合意形成"
            ]),
            ("中期（3-6ヶ月）", [
                "優先施策の実行開始",
                "KPIの設定とモニタリング体制の構築",
                "中間評価と軌道修正"
            ]),
            ("長期（6-12ヶ月）", [
                "施策の効果測定と評価",
                "次フェーズの戦略立案",
                "継続的改善サイクルの確立"
            ])
        ]
        
        for period, actions in action_plans:
            # 期間
            p = text_frame.add_paragraph()
            p.text = period
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            p.space_after = Pt(6)
            
            # アクション
            for action in actions:
                p = text_frame.add_paragraph()
                p.text = f"  • {action}"
                p.level = 1
                p.font.size = Pt(14)
                p.font.color.rgb = self.colors['text']
            
            # 空行
            text_frame.add_paragraph()
